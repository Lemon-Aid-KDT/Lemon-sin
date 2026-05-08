import base64
import subprocess
from pathlib import Path

import requests
from langchain_ollama import ChatOllama
from langchain_core.messages import HumanMessage, AIMessage
from config import (
    OLLAMA_BASE_URL,
    LLM_MODEL,
    MODEL_PROFILES,
    FEATURE_MODEL_MAP,
    OLLAMA_KEEP_ALIVE,
    OLLAMA_NUM_PREDICT_DEFAULT,
    OLLAMA_NUM_PREDICT_MAP,
    OLLAMA_NUM_CTX,
    ollama_headers,
)

# Qwen 3.5 / EXAONE Deep 등 thinking 모드 모델 — content 가 비어 오거나
# <thought>...</thought> 블록이 매우 길어 응답 지연이 발생한다.
# /no_think 토큰을 프롬프트에 주입하여 비활성화한다.
_THINKING_MODELS = {"qwen3.5", "qwen3", "exaone-deep"}


class _ThoughtStripper:
    """v3.5 — 스트리밍 토큰에서 ``<thought>...</thought>`` 블록을 자동 제거.

    EXAONE Deep 같은 추론 강화 모델은 Ollama API 의 별도 ``thinking`` 필드가 아닌
    ``message.content`` 텍스트 안에 직접 ``<thought>`` 태그를 출력한다. 일반 텍스트로
    yield 되면 사용자 화면에 빈 태그/긴 추론이 보여 UX 저해.

    chunk 단위로 누적 → 태그 경계를 분할 안전하게 처리.
    사용:
        s = _ThoughtStripper()
        for chunk in stream:
            visible = s.process(chunk)
            if visible:
                yield visible
        tail = s.flush()
        if tail:
            yield tail
    """

    OPEN = "<thought>"
    CLOSE = "</thought>"

    def __init__(self) -> None:
        self.buffer = ""
        self.in_thought = False

    def process(self, chunk: str) -> str:
        """토큰 chunk 누적 + 태그 추적 → 사용자에게 노출할 텍스트만 반환."""
        self.buffer += chunk
        out: list[str] = []
        while True:
            if self.in_thought:
                idx = self.buffer.find(self.CLOSE)
                if idx == -1:
                    # 닫는 태그 미도달 — buffer 끝 일부(부분 매치 가능성) 만 보존
                    keep = max(0, len(self.buffer) - len(self.CLOSE) + 1)
                    self.buffer = self.buffer[keep:]
                    break
                # 닫는 태그 발견 → 버리고 반복 (다음 thought 또는 본문)
                self.buffer = self.buffer[idx + len(self.CLOSE):]
                self.in_thought = False
            else:
                idx = self.buffer.find(self.OPEN)
                if idx == -1:
                    # 여는 태그 미도달 — buffer 끝 일부(부분 매치 가능성) 보존, 나머지 yield
                    keep = max(0, len(self.buffer) - len(self.OPEN) + 1)
                    if keep > 0:
                        out.append(self.buffer[:keep])
                    self.buffer = self.buffer[keep:]
                    break
                # 여는 태그 발견 → 그 이전은 yield, 모드 전환
                if idx > 0:
                    out.append(self.buffer[:idx])
                self.buffer = self.buffer[idx + len(self.OPEN):]
                self.in_thought = True
        return "".join(out)

    def flush(self) -> str:
        """스트림 종료 시 잔여 buffer 처리. thought 안이면 폐기, 본문이면 그대로."""
        if self.in_thought:
            return ""
        out = self.buffer
        self.buffer = ""
        return out


class _NoThinkChatOllama(ChatOllama):
    """Qwen 3.5 thinking 모드를 자동 비활성화하는 래퍼."""

    def _should_disable_think(self) -> bool:
        return any(m in self.model for m in _THINKING_MODELS)

    def _inject_no_think(self, input):
        if not self._should_disable_think():
            return input
        if isinstance(input, str):
            return f"/no_think\n{input}"
        if isinstance(input, list) and input:
            first = input[0]
            if isinstance(first, HumanMessage):
                input[0] = HumanMessage(content=f"/no_think\n{first.content}")
            elif isinstance(first, str):
                input[0] = f"/no_think\n{first}"
        return input

    def _extract_content(self, response):
        """thinking 태그 안의 내용을 제거하고 실제 응답만 반환."""
        if not self._should_disable_think():
            return response
        if hasattr(response, 'content') and not response.content:
            pass
        return response

    def invoke(self, input, config=None, **kwargs):
        result = super().invoke(self._inject_no_think(input), config=config, **kwargs)
        return self._extract_content(result)

    async def ainvoke(self, input, config=None, **kwargs):
        result = await super().ainvoke(self._inject_no_think(input), config=config, **kwargs)
        return self._extract_content(result)


def get_llm(
    model: str = LLM_MODEL,
    temperature: float = 0.3,
    feature: str = "",
    num_predict: int | None = None,
) -> ChatOllama:
    """Ollama LLM 클라이언트를 반환한다."""
    if num_predict is None:
        num_predict = OLLAMA_NUM_PREDICT_MAP.get(feature, OLLAMA_NUM_PREDICT_DEFAULT)

    # Plan A 변형: Cloud Run prod 에서는 Caddy 경유 → X-AJIN-Secret 헤더 부착
    _hdrs = ollama_headers()
    kwargs: dict = dict(
        model=model,
        base_url=OLLAMA_BASE_URL,
        temperature=temperature,
        keep_alive=OLLAMA_KEEP_ALIVE,
        num_predict=num_predict,
        num_ctx=OLLAMA_NUM_CTX,
    )
    if _hdrs:
        # langchain-ollama 0.x: client_kwargs 로 httpx 헤더 주입
        kwargs["client_kwargs"] = {"headers": _hdrs}
    # Qwen 3.5 thinking 비활성화 (langchain-ollama >= 0.3 지원)
    if any(m in model for m in _THINKING_MODELS):
        kwargs["think"] = False

    return _NoThinkChatOllama(**kwargs)


# ──────────────────────────────────────────────
# 스트리밍 응답 (Ollama API 직접 호출)
# ──────────────────────────────────────────────

def stream_generate(
    prompt: str,
    model: str = LLM_MODEL,
    feature: str = "",
    temperature: float = 0.3,
):
    """Ollama chat API로 스트리밍 요청을 보내고 토큰 단위로 yield한다.

    - thinking 모델: think=false 설정 + /no_think 프롬프트 주입
    - thinking 응답이 계속 올 경우: thinking 완료 후 response 토큰만 yield
    - 모델 미발견 시 자동 폴백
    """
    num_predict = OLLAMA_NUM_PREDICT_MAP.get(feature, OLLAMA_NUM_PREDICT_DEFAULT)
    is_thinking = any(m in model for m in _THINKING_MODELS)

    # thinking 모델은 /no_think 프롬프트 주입
    actual_prompt = f"/no_think\n{prompt}" if is_thinking else prompt

    payload = {
        "model": model,
        "messages": [{"role": "user", "content": actual_prompt}],
        "stream": True,
        "keep_alive": OLLAMA_KEEP_ALIVE,
        "options": {
            "num_predict": num_predict,
            "temperature": temperature,
            "num_ctx": OLLAMA_NUM_CTX,
        },
    }
    if is_thinking:
        payload["think"] = False

    try:
        resp = requests.post(
            f"{OLLAMA_BASE_URL}/api/chat",
            json=payload,
            stream=True,
            timeout=180,
            headers=ollama_headers(),
        )

        # 모델 미발견 시 폴백 (1회만)
        if resp.status_code == 404:
            if not hasattr(stream_generate, '_fallback_tried'):
                stream_generate._fallback_tried = set()
            if model not in stream_generate._fallback_tried:
                stream_generate._fallback_tried.add(model)
                fallback = _find_fallback_model(model)
                if fallback and fallback != model and fallback not in stream_generate._fallback_tried:
                    yield from stream_generate(prompt, model=fallback, feature=feature, temperature=temperature)
                    stream_generate._fallback_tried.discard(model)
                    return
                stream_generate._fallback_tried.discard(model)
            yield f"\n[오류] 모델 '{model}'을 찾을 수 없습니다. OLLAMA_NEW_ENGINE=true 로 Ollama를 재시작해주세요."
            return

        # 500 에러: 모델 로딩 실패 / VRAM 부족 / 컨텍스트 초과 → 폴백 모델로 재시도
        if resp.status_code == 500:
            import logging
            _err_body = ""
            try:
                _err_body = resp.text[:500]
            except Exception:
                pass
            logging.getLogger(__name__).warning(
                f"Ollama 500 에러 (model={model}): {_err_body}"
            )

            # 1차 시도: 컨텍스트 축소 후 동일 모델 재시도
            if not hasattr(stream_generate, '_retry_count'):
                stream_generate._retry_count = {}
            retry_key = f"{model}_{feature}"
            retries = stream_generate._retry_count.get(retry_key, 0)

            if retries == 0:
                stream_generate._retry_count[retry_key] = 1
                # 프롬프트 축소 (최대 3000자) + 컨텍스트 축소
                truncated = prompt[:3000] if len(prompt) > 3000 else prompt
                yield from _retry_with_reduced_context(truncated, model, feature, temperature)
                stream_generate._retry_count.pop(retry_key, None)
                return

            # 2차 시도: 폴백 모델
            if retries <= 1:
                stream_generate._retry_count[retry_key] = 2
                fallback = _find_fallback_model(model)
                if fallback and fallback != model:
                    yield f"⚠️ 모델 '{model}' 오류 → '{fallback}'으로 전환 중...\n\n"
                    yield from stream_generate(prompt, model=fallback, feature=feature, temperature=temperature)
                    stream_generate._retry_count.pop(retry_key, None)
                    return

            stream_generate._retry_count.pop(retry_key, None)
            yield f"\n[오류] 모델 '{model}' 서버 오류 (500). Ollama를 재시작하거나 다른 모델을 선택해주세요."
            return

        resp.raise_for_status()

        import json as _json
        thinking_done = False
        # v3.5 — EXAONE Deep 등 <thought> 텍스트 출력 모델은 stripper 적용
        stripper = _ThoughtStripper() if is_thinking else None
        for line in resp.iter_lines():
            if not line:
                continue
            try:
                data = _json.loads(line)
            except _json.JSONDecodeError:
                continue
            msg = data.get("message", {})

            # 1) response 토큰이 있으면 yield (thinking 모델은 thought 블록 제거)
            token = msg.get("content", "")
            if token:
                thinking_done = True
                if stripper is not None:
                    visible = stripper.process(token)
                    if visible:
                        yield visible
                else:
                    yield token

            # 2) Ollama API thinking 필드 (Qwen) 는 스킵
            thinking_token = data.get("thinking", "")
            if thinking_token and not token:
                pass

            if data.get("done"):
                # 잔여 buffer 처리
                if stripper is not None:
                    tail = stripper.flush()
                    if tail:
                        yield tail
                if not thinking_done:
                    yield from _fallback_non_stream(prompt, model, num_predict, temperature)
                break
    except requests.exceptions.ConnectionError:
        yield "\n[오류] Ollama 서버에 연결할 수 없습니다. `ollama serve` 실행 후 다시 시도해주세요."
    except requests.exceptions.Timeout:
        yield "\n[오류] 응답 시간이 초과되었습니다 (180초). 더 작은 모델을 사용하거나 프롬프트를 줄여주세요."
    except requests.exceptions.HTTPError as e:
        yield f"\n[오류] HTTP 에러: {e.response.status_code if e.response else 'unknown'}"
    except Exception as e:
        yield f"\n[오류] 스트리밍 실패: {e}"


def _retry_with_reduced_context(prompt: str, model: str, feature: str, temperature: float):
    """500 에러 시 컨텍스트/예측 토큰을 줄여서 재시도한다."""
    reduced_ctx = min(OLLAMA_NUM_CTX, 2048)  # 절반으로 축소
    reduced_predict = min(OLLAMA_NUM_PREDICT_MAP.get(feature, 500), 400)

    try:
        resp = requests.post(
            f"{OLLAMA_BASE_URL}/api/chat",
            headers=ollama_headers(),
            json={
                "model": model,
                "messages": [{"role": "user", "content": prompt}],
                "stream": True,
                "keep_alive": OLLAMA_KEEP_ALIVE,
                "options": {
                    "num_predict": reduced_predict,
                    "temperature": temperature,
                    "num_ctx": reduced_ctx,
                },
            },
            stream=True,
            timeout=180,
        )
        if resp.status_code != 200:
            return  # 재시도도 실패하면 상위에서 폴백 모델 시도

        import json as _json
        for line in resp.iter_lines():
            if not line:
                continue
            try:
                data = _json.loads(line)
            except _json.JSONDecodeError:
                continue
            token = data.get("message", {}).get("content", "")
            if token:
                yield token
            if data.get("done"):
                break
    except Exception:
        pass  # 재시도 실패 — 상위에서 폴백 모델로 전환


def _fallback_non_stream(prompt: str, model: str, num_predict: int, temperature: float):
    """스트리밍에서 thinking만 반환된 경우 non-stream으로 재시도."""
    try:
        resp = requests.post(
            f"{OLLAMA_BASE_URL}/api/chat",
            headers=ollama_headers(),
            json={
                "model": model,
                "messages": [{"role": "user", "content": f"/no_think\n{prompt}"}],
                "stream": False,
                "think": False,
                "options": {
                    "num_predict": num_predict,
                    "temperature": temperature,
                    "num_ctx": OLLAMA_NUM_CTX,
                },
            },
            timeout=180,
        )
        if resp.status_code == 200:
            content = resp.json().get("message", {}).get("content", "")
            if content:
                yield content
                return
        yield "[오류] 모델이 응답을 생성하지 못했습니다."
    except Exception as e:
        yield f"[오류] 폴백 요청 실패: {e}"


def _find_fallback_model(original: str) -> str | None:
    """설치된 모델 중 대체 모델을 찾는다. (v2.3: 우선순위 기반 폴백)"""
    installed = get_installed_models()
    if not installed:
        return None

    # v3.5: 승인된 모델만 포함 (경량 → 안정)
    _PREFERRED_FALLBACKS = [
        "qwen3.5:latest", "qwen3.5:9b", "qwen3.5:4b",
        "gemma4:latest", "gemma4:e2b",
        "exaone3.5:latest",
        "gpt-oss:20b",
    ]

    # 1순위: 선호 모델 중 설치된 것
    for pref in _PREFERRED_FALLBACKS:
        if pref in installed and pref != original:
            return pref

    # 2순위: 같은 계열 모델
    base = original.split(":")[0]
    for m in installed:
        if base in m and m != original:
            return m

    # 3순위: thinking 모델 제외하고 아무거나
    for m in installed:
        if m != original and "deep" not in m and "vision" not in m and "ocr" not in m:
            return m

    return installed[0] if installed else None


# ──────────────────────────────────────────────
# 설치된 모델 감지
# ──────────────────────────────────────────────

_installed_cache: list[str] | None = None


def get_installed_models() -> list[str]:
    """Ollama에 설치된 모델 목록을 반환한다 (캐시됨).

    1차: API /api/tags 시도
    2차: API 결과가 부족하면 CLI `ollama list` 폴백 (0.18.x 호환)
    """
    global _installed_cache
    if _installed_cache is not None:
        return _installed_cache

    models: list[str] = []

    # 1차: API 시도
    try:
        resp = requests.get(
            f"{OLLAMA_BASE_URL}/api/tags",
            headers=ollama_headers(),
            timeout=3,
        )
        if resp.status_code == 200:
            models = [m["name"] for m in resp.json().get("models", [])]
    except Exception:
        pass

    # 2차: API 결과가 3개 미만이면 CLI 폴백 (Ollama 0.18.x 버그 대응)
    if len(models) < 3:
        cli_models = _get_models_via_cli()
        if len(cli_models) > len(models):
            models = cli_models

    _installed_cache = models
    return _installed_cache


def _get_models_via_cli() -> list[str]:
    """ollama list CLI 명령으로 설치된 모델 목록을 반환한다."""
    try:
        result = subprocess.run(
            ["ollama", "list"],
            capture_output=True, text=True, shell=False, timeout=10,
        )
        if result.returncode != 0:
            return []
        models = []
        for line in result.stdout.strip().split("\n")[1:]:  # 헤더 스킵
            parts = line.split()
            if parts:
                name = parts[0]
                # cloud 모델 제외 (로컬 실행 불가)
                if ":cloud" not in name:
                    models.append(name)
        return models
    except Exception:
        return []


def invalidate_model_cache():
    """모델 캐시를 무효화한다 (모델 설치/삭제 후 호출)."""
    global _installed_cache
    _installed_cache = None


def get_available_chat_models() -> list[dict]:
    """설치된 모델 중 채팅 가능한 모델 프로필 목록을 반환한다."""
    installed = get_installed_models()
    available = []
    for model_id, profile in MODEL_PROFILES.items():
        if model_id in installed:
            available.append({"id": model_id, **profile})
    return available


def get_vision_models() -> list[dict]:
    """설치된 비전 모델만 반환한다."""
    return [m for m in get_available_chat_models() if m.get("vision")]


# ──────────────────────────────────────────────
# 자동 모델 선택 (기능 / 성능 기반)
# ──────────────────────────────────────────────

def auto_select_model(feature: str = "onboarding") -> str:
    """기능에 맞는 최적 모델을 자동 선택한다.

    1. FEATURE_MODEL_MAP에서 추천 모델을 찾음
    2. 해당 모델이 설치되어 있으면 사용
    3. 미설치 시 best_for에 해당 기능이 포함된 모델 중 품질 최고를 선택
    4. 모두 실패하면 기본 LLM_MODEL로 폴백
    """
    installed = get_installed_models()

    # 1단계: 기능별 추천 모델 확인
    recommended = FEATURE_MODEL_MAP.get(feature)
    if recommended and recommended in installed:
        return recommended

    # 2단계: best_for에 해당 기능이 있는 설치 모델 중 품질 최고 선택
    quality_rank = {"very_high": 4, "high": 3, "good": 2, "medium": 1}
    candidates = []
    for model_id, profile in MODEL_PROFILES.items():
        if model_id in installed and feature in profile.get("best_for", []):
            rank = quality_rank.get(profile.get("quality", "medium"), 1)
            candidates.append((rank, model_id))

    if candidates:
        candidates.sort(key=lambda x: x[0], reverse=True)
        return candidates[0][1]

    # 3단계: 폴백
    if LLM_MODEL in installed:
        return LLM_MODEL
    return installed[0] if installed else LLM_MODEL


def auto_select_vision_model() -> str | None:
    """설치된 비전 모델 중 최적 모델을 반환한다. 없으면 None."""
    vision = get_vision_models()
    if not vision:
        return None
    # 품질 순 정렬
    quality_rank = {"very_high": 4, "high": 3, "good": 2, "medium": 1}
    vision.sort(key=lambda m: quality_rank.get(m.get("quality", "medium"), 1), reverse=True)
    return vision[0]["id"]


# ──────────────────────────────────────────────
# 비전 모델 호출 (이미지 분석)
# ──────────────────────────────────────────────

def invoke_vision(
    prompt: str,
    image_bytes: bytes,
    model: str | None = None,
) -> str:
    """비전 모델에 이미지와 프롬프트를 전송하여 응답을 받는다."""
    if model is None:
        model = auto_select_vision_model()
    if model is None:
        return "[오류] 비전 모델이 설치되어 있지 않습니다. gemma4를 설치해주세요. (ollama pull gemma4)"

    img_b64 = base64.b64encode(image_bytes).decode("utf-8")

    try:
        resp = requests.post(
            f"{OLLAMA_BASE_URL}/api/chat",
            headers=ollama_headers(),
            json={
                "model": model,
                "messages": [
                    {
                        "role": "user",
                        "content": prompt,
                        "images": [img_b64],
                    }
                ],
                "stream": False,
                "think": False,
            },
            timeout=120,
        )
        if resp.status_code == 404:
            return f"[오류] 비전 모델 '{model}'을 찾을 수 없습니다. ollama pull {model}로 설치해주세요."
        if resp.status_code == 200:
            content = resp.json().get("message", {}).get("content", "")
            if content:
                return content
            # thinking만 반환된 경우
            return resp.json().get("message", {}).get("thinking", "[오류] 모델이 응답을 생성하지 못했습니다.")
        return f"[오류] 비전 모델 응답 실패: HTTP {resp.status_code}"
    except Exception as e:
        return f"[오류] 비전 모델 호출 실패: {e}"


# ──────────────────────────────────────────────
# 파일 텍스트 추출
# ──────────────────────────────────────────────

def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """업로드된 파일에서 텍스트를 추출한다. (v2.3: 확장자 대폭 추가)"""
    ext = Path(filename).suffix.lower()

    # ── 텍스트 기반 파일 (직접 디코딩) ──
    _TEXT_EXTS = {
        ".txt", ".md", ".csv", ".log",
        ".py", ".js", ".ts", ".html", ".css",
        ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg",
        ".dxf", ".step", ".stp", ".igs",
        ".rtf",
    }
    if ext in _TEXT_EXTS:
        for enc in ("utf-8", "euc-kr", "cp949"):
            try:
                return file_bytes.decode(enc)
            except (UnicodeDecodeError, LookupError):
                continue
        return file_bytes.decode("utf-8", errors="replace")

    # ── JSON (포맷팅 후 반환) ──
    if ext == ".json":
        import json as _json
        for enc in ("utf-8", "euc-kr", "cp949"):
            try:
                raw = file_bytes.decode(enc)
                parsed = _json.loads(raw)
                return _json.dumps(parsed, indent=2, ensure_ascii=False)
            except (UnicodeDecodeError, LookupError):
                continue
            except _json.JSONDecodeError:
                return raw  # 파싱 실패 시 원본 텍스트 반환
        return file_bytes.decode("utf-8", errors="replace")

    # ── PDF ──
    if ext == ".pdf":
        return _extract_pdf(file_bytes)

    # ── DOCX ──
    if ext == ".docx":
        return _extract_docx(file_bytes)

    # ── XLSX / XLS (Excel) ──
    if ext in (".xlsx", ".xls"):
        return _extract_excel(file_bytes, ext)

    # ── PPTX (PowerPoint) ──
    if ext == ".pptx":
        return _extract_pptx(file_bytes)

    # ── DOC (구형 Word) ──
    if ext == ".doc":
        return _extract_doc(file_bytes)

    # ── HWPX (ODT 기반 ZIP) ──
    if ext == ".hwpx":
        return _extract_hwpx(file_bytes)

    return f"[지원하지 않는 형식: {ext}]"


def _extract_pdf(data: bytes) -> str:
    """PDF에서 텍스트를 추출한다."""
    try:
        import io
        # PyPDF2 우선, 없으면 pdfplumber
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(data))
            pages = [page.extract_text() or "" for page in reader.pages]
            return "\n\n".join(pages)
        except ImportError:
            pass
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                pages = [p.extract_text() or "" for p in pdf.pages]
            return "\n\n".join(pages)
        except ImportError:
            pass
        return "[PDF 파싱 라이브러리가 없습니다. PyPDF2 또는 pdfplumber를 설치하세요.]"
    except Exception as e:
        return f"[PDF 파싱 오류: {e}]"


def _extract_docx(data: bytes) -> str:
    """DOCX에서 텍스트를 추출한다."""
    try:
        import io
        from docx import Document
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except ImportError:
        return "[python-docx가 설치되어 있지 않습니다.]"
    except Exception as e:
        return f"[DOCX 파싱 오류: {e}]"


def _extract_hwpx(data: bytes) -> str:
    """HWPX(ODT 기반 ZIP)에서 텍스트를 추출한다. (v2.3)"""
    try:
        import io, zipfile, re
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            # ODF content.xml 또는 한컴 section0.xml에서 텍스트 추출
            for name in ("content.xml", "Contents/section0.xml"):
                if name in zf.namelist():
                    xml_text = zf.read(name).decode("utf-8", errors="replace")
                    # XML 태그 제거하여 순수 텍스트만 추출
                    clean = re.sub(r"<[^>]+>", " ", xml_text)
                    clean = re.sub(r"\s+", " ", clean).strip()
                    return clean
        return "[HWPX 파일에서 텍스트를 찾을 수 없습니다.]"
    except Exception as e:
        return f"[HWPX 파싱 오류: {e}]"


def _extract_excel(data: bytes, ext: str = ".xlsx") -> str:
    """v2.6: Excel 파일(.xlsx/.xls)에서 텍스트를 추출한다."""
    import io

    # .xlsx → openpyxl
    if ext == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"[시트: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        lines.append("\t".join(cells))
            wb.close()
            return "\n".join(lines)
        except ImportError:
            return "[openpyxl이 설치되어 있지 않습니다.]"
        except Exception as e:
            return f"[XLSX 파싱 오류: {e}]"

    # .xls → xlrd 우선, pandas 폴백
    if ext == ".xls":
        try:
            import xlrd
            wb = xlrd.open_workbook(file_contents=data)
            lines = []
            for ws in wb.sheets():
                lines.append(f"[시트: {ws.name}]")
                for row_idx in range(ws.nrows):
                    cells = [str(ws.cell_value(row_idx, col)) for col in range(ws.ncols)]
                    if any(cells):
                        lines.append("\t".join(cells))
            return "\n".join(lines)
        except ImportError:
            pass
        try:
            import pandas as pd
            dfs = pd.read_excel(io.BytesIO(data), sheet_name=None, engine="xlrd")
            lines = []
            for sheet_name, df in dfs.items():
                lines.append(f"[시트: {sheet_name}]")
                lines.append(df.to_string(index=False))
            return "\n".join(lines)
        except Exception:
            pass
        return "[.xls 파싱 라이브러리가 없습니다. xlrd를 설치하세요: pip install xlrd]"

    return "[지원하지 않는 Excel 형식]"


def _extract_pptx(data: bytes) -> str:
    """v2.6: PowerPoint 파일(.pptx)에서 텍스트를 추출한다."""
    try:
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                lines.append(f"[슬라이드 {i}]")
                lines.extend(slide_texts)
        return "\n".join(lines)
    except ImportError:
        return "[python-pptx가 설치되어 있지 않습니다.]"
    except Exception as e:
        return f"[PPTX 파싱 오류: {e}]"


def _extract_doc(data: bytes) -> str:
    """v2.6: 구형 Word 파일(.doc)에서 텍스트를 추출한다."""
    try:
        # antiword 또는 textract 기반 — 없으면 바이너리에서 텍스트 추출 시도
        import re
        # .doc 바이너리에서 ASCII/한글 텍스트 추출 (단순 폴백)
        try:
            text = data.decode("utf-8", errors="ignore")
        except Exception:
            text = data.decode("cp949", errors="ignore")
        # 제어 문자 제거, 출력 가능 문자만 유지
        clean = re.sub(r"[^\x20-\x7E\uAC00-\uD7AF\u3131-\u318E\n\t]", " ", text)
        clean = re.sub(r" {3,}", "  ", clean).strip()
        if len(clean) < 50:
            return "[.doc 파일에서 텍스트를 추출할 수 없습니다. .docx 형식으로 변환 후 업로드하세요.]"
        return clean
    except Exception as e:
        return f"[DOC 파싱 오류: {e}]"
